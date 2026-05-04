from flask import Flask, jsonify, render_template, request, flash, session, send_file, redirect, url_for
from groq import Groq
from flask_sqlalchemy import SQLAlchemy
from werkzeug.security import generate_password_hash, check_password_hash
import google.generativeai as genai
from dotenv import load_dotenv
import os
import re
import io
import base64
from datetime import datetime
import requests
# ORIGINAL SPECIALIZED IMPORTS
import img2pdf
from PyPDF2 import PdfReader
import yt_dlp

# ADDED FOR POWERPOINT GENERATION
from pptx import Presentation
from pptx.util import Inches

load_dotenv()

app = Flask(__name__)
app.secret_key = "smartstudy-secret-key-2026"

# ==================== DATABASE & AUTH CONFIGURATION ====================
# Absolute path for PythonAnywhere stability
BASE_DIR = os.path.abspath(os.path.dirname(__file__))
app.config['SQLALCHEMY_DATABASE_URI'] = 'sqlite:///' + os.path.join(BASE_DIR, 'users.db')
app.config['SQLALCHEMY_TRACK_MODIFICATIONS'] = False
db = SQLAlchemy(app)

class User(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    name = db.Column(db.String(100), nullable=False)
    email = db.Column(db.String(120), unique=True, nullable=False)
    password = db.Column(db.String(200), nullable=False)
    subjects = db.relationship('Subject', backref='user', lazy=True)

class Subject(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    name = db.Column(db.String(100), nullable=False)
    user_id = db.Column(db.Integer, db.ForeignKey('user.id'), nullable=False)
    lectures = db.relationship('Lecture', backref='subject', cascade="all, delete-orphan", lazy=True)

class Lecture(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    topic = db.Column(db.String(200), nullable=False)
    sub_topics = db.Column(db.Text)
    ai_concepts = db.Column(db.Text)
    ai_quiz = db.Column(db.Text)
    date_added = db.Column(db.DateTime, default=datetime.utcnow)
    subject_id = db.Column(db.Integer, db.ForeignKey('subject.id'), nullable=False)

with app.app_context():
    db.create_all()

# ==================== CLOUD API CONFIGURATION ====================
GENAI_API_KEY = "Upload your own" # PASTE YOUR ACTUAL KEY HERE

DOWNLOADS_DIR = "downloads"
os.makedirs(DOWNLOADS_DIR, exist_ok=True)
os.makedirs("uploads", exist_ok=True)



# ==================== CORE LLM FUNCTION (GEMINI) ====================

# Replace with your Groq API Key
client = Groq(api_key="Upload your own")

def ask_local_llm(prompt, model_type="fast", image_data=None):
    try:
        # Groq uses Llama 3 models which are incredibly fast for quizzes/summaries
        completion = client.chat.completions.create(
            model="llama-3.3-70b-versatile",
            messages=[{"role": "user", "content": prompt}],
        )
        return completion.choices[0].message.content
    except Exception as e:
        return f"Groq Error: {str(e)}"

        # ...end of API configuration function....

def chunk_text(text, size=1000):
    return [text[i:i+size] for i in range(0, len(text), size)]

# ==================== AUTH ROUTES ====================

@app.route('/signup', methods=['GET', 'POST'])
def signup():
    if request.method == 'POST':
        name = request.form.get('name')
        email = request.form.get('email')
        password = request.form.get('password')
        user_exists = User.query.filter_by(email=email).first()
        if user_exists:
            flash("Email already exists!", "danger")
        else:
            hashed_pw = generate_password_hash(password, method='pbkdf2:sha256')
            new_user = User(name=name, email=email, password=hashed_pw)
            db.session.add(new_user)
            db.session.commit()
            flash("Signup successful! Please login.", "success")
            return redirect(url_for('login'))
    return render_template('signup.html')

@app.route('/login', methods=['GET', 'POST'])
def login():
    if request.method == 'POST':
        email = request.form.get('email')
        password = request.form.get('password')
        user = User.query.filter_by(email=email).first()
        if user and check_password_hash(user.password, password):
            session['user_id'] = user.id
            session['user_name'] = user.name
            session['user_email'] = user.email
            flash(f"Welcome back, {user.name}!", "success")
            return redirect(url_for('index'))
        else:
            flash("Login failed. Check your email and password.", "danger")
    return render_template('login.html')

@app.route('/logout')
def logout():
    session.clear()
    flash("You have been logged out.", "info")
    return redirect(url_for('index'))

# ==================== COURSE COMPENDIUM ROUTES ====================

@app.route('/courses', methods=['GET', 'POST'])
def courses():
    if 'user_id' not in session:
        flash("Please login to use Course Compendium", "warning")
        return redirect(url_for('login'))
    if request.method == 'POST':
        subject_name = request.form.get('subject_name')
        if subject_name:
            new_sub = Subject(name=subject_name, user_id=session['user_id'])
            db.session.add(new_sub)
            db.session.commit()
            flash(f"Subject '{subject_name}' added!", "success")
    user_subjects = Subject.query.filter_by(user_id=session['user_id']).all()
    return render_template('courses.html', subjects=user_subjects)

@app.route('/rename_subject/<int:subject_id>', methods=['POST'])
def rename_subject(subject_id):
    if 'user_id' not in session: return redirect(url_for('login'))
    subject = Subject.query.get_or_404(subject_id)
    if subject.user_id != session['user_id']: return redirect(url_for('courses'))
    new_name = request.form.get('new_name')
    if new_name:
        subject.name = new_name
        db.session.commit()
        flash("Folder renamed!", "success")
    return redirect(url_for('courses'))

@app.route('/delete_subject/<int:subject_id>')
def delete_subject(subject_id):
    if 'user_id' not in session: return redirect(url_for('login'))
    subject = Subject.query.get_or_404(subject_id)
    if subject.user_id != session['user_id']: return redirect(url_for('courses'))
    db.session.delete(subject)
    db.session.commit()
    return redirect(url_for('courses'))

@app.route('/delete_lecture/<int:lecture_id>')
def delete_lecture(lecture_id):
    if 'user_id' not in session: return redirect(url_for('login'))
    lecture = Lecture.query.get_or_404(lecture_id)
    subject_id = lecture.subject_id
    if lecture.subject.user_id != session['user_id']: return redirect(url_for('courses'))
    db.session.delete(lecture)
    db.session.commit()
    return redirect(url_for('view_subject', subject_id=subject_id))

@app.route('/subject/<int:subject_id>', methods=['GET', 'POST'])
def view_subject(subject_id):
    subject = Subject.query.get_or_404(subject_id)
    if request.method == 'POST':
        topic = request.form.get('topic')
        sub_topics = request.form.get('sub_topics')
        action = request.form.get('action')
        concepts, quiz = "", ""
        if action == 'save_gen':
            concepts = ask_local_llm(f"Extract core concepts for {topic}: {sub_topics}")
            quiz = ask_local_llm(f"Generate 3 MCQs with answers for {topic}: {sub_topics}")
        new_lecture = Lecture(topic=topic, sub_topics=sub_topics, ai_concepts=concepts, ai_quiz=quiz, subject_id=subject_id)
        db.session.add(new_lecture)
        db.session.commit()
        flash("Work logged!", "success")
    return render_template('view_subject.html', subject=subject)

@app.route('/generate_kit/<int:lecture_id>')
def generate_kit(lecture_id):
    lecture = Lecture.query.get_or_404(lecture_id)
    lecture.ai_concepts = ask_local_llm(f"Extract concepts for {lecture.topic}: {lecture.sub_topics}")
    lecture.ai_quiz = ask_local_llm(f"Generate 3 MCQs for {lecture.topic}: {lecture.sub_topics}")
    db.session.commit()
    return redirect(url_for('view_subject', subject_id=lecture.subject_id))

# ==================== TOOL ROUTES ====================

@app.route('/')
def index():
    return render_template('index.html', user_name=session.get('user_name'))

@app.route('/diagram')
def diagram_studio():
    return render_template('diagram.html')

@app.route('/diagram-ai', methods=['POST'])
def diagram_ai():
    data = request.json
    user_prompt = data.get('prompt', '')
    sys_instr = "Generate ONLY Mermaid.js code starting with 'graph TD'. Rules: Use pipe format A -->|Label| B. Node IDs single letters. Brackets for text A[Text]. No colons."
    mermaid_code = ask_local_llm(f"{sys_instr}\n\nUser: {user_prompt}")
    cleaned = re.sub(r'```mermaid|```', '', mermaid_code).strip()
    return jsonify({"code": cleaned})

@app.route('/diagram-sketch', methods=['POST'])
def diagram_sketch():
    file = request.files.get('sketch')
    if not file: return jsonify({"error": "No file"}), 400
    image_data = base64.b64encode(file.read()).decode('utf-8')
    prompt = "Convert this sketch into Mermaid.js 'graph TD' code. Output only code."
    mermaid_code = ask_local_llm(prompt, image_data=image_data)
    cleaned = re.sub(r'```mermaid|```', '', mermaid_code).strip()
    return jsonify({"code": cleaned})

@app.route('/get_response', methods=['POST'])
def get_response():
    user_name = session.get('user_name', 'Guest')
    user_msg = request.json.get('message') if request.is_json else request.form.get('message')
    image_file = request.files.get('image')
    base64_image = base64.b64encode(image_file.read()).decode('utf-8') if image_file else None

    context = session.get('last_summary', '')
    prompt = f"User: {user_name}. Context: {context}\nQuestion: {user_msg}"
    reply = ask_local_llm(prompt, image_data=base64_image)
    return jsonify({"response": reply})

@app.route('/notes', methods=['GET', 'POST'])
def notes():
    answer, question = None, None
    # Path for a temporary context file
    context_file = os.path.join(BASE_DIR, f"context_{session.get('user_id', 'anon')}.txt")

    if request.method == 'POST':
        if 'pdf' in request.files:
            file = request.files['pdf']
            if file and file.filename.endswith('.pdf'):
                reader = PdfReader(file)
                text = "".join([p.extract_text() or "" for p in reader.pages])

                # SAVE TO FILE INSTEAD OF SESSION
                with open(context_file, "w", encoding="utf-8") as f:
                    f.write(text[:20000]) # You can store more now!

                session['pdf_indexed'] = True # Just a small flag
                flash("PDF Indexed successfully!", "success")

        elif 'question' in request.form:
            question = request.form.get('question')

            # READ FROM FILE
            context = ""
            if os.path.exists(context_file):
                with open(context_file, "r", encoding="utf-8") as f:
                    context = f.read()

            prompt = f"Using this text: {context}\n\nAnswer: {question}"
            answer = ask_local_llm(prompt)

    return render_template('notes.html', answer=answer, question=question)

@app.route('/images-to-pdf', methods=['GET', 'POST'])
def images_to_pdf():
    if request.method == 'POST':
        files = request.files.getlist("images")
        fmt = request.form.get("format")
        img_bytes_list = [f.read() for f in files]
        if fmt == 'ppt':
            prs = Presentation()
            for img in img_bytes_list:
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                prs.slides[prs.slides.index(slide)].shapes.add_picture(io.BytesIO(img), 0, 0, width=Inches(13.33))
            out = io.BytesIO(); prs.save(out); out.seek(0)
            return send_file(out, mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation', as_attachment=True, download_name='Slides.pptx')
        else:
            pdf = img2pdf.convert(img_bytes_list)
            return send_file(io.BytesIO(pdf), mimetype='application/pdf', as_attachment=True, download_name='Notes.pdf')
    return render_template('images_to_pdf.html')

@app.route('/planner', methods=['GET', 'POST'])
def planner():
    plan = None
    if request.method == 'POST':
        subjects = request.form.get('subjects', '')
        hours = request.form.get('hours', '4')
        days = request.form.get('days', '7')
        prompt = f"Create a study plan for: {subjects} for {hours}h/day over {days} days."
        plan = ask_local_llm(prompt)
    return render_template('planner.html', plan=plan)

@app.route('/quiz', methods=['GET', 'POST'])
def quiz():
    quiz_content = None
    if request.method == 'POST':
        topic = request.form.get('topic', '')
        num = request.form.get('num', '5')
        prompt = f"Generate {num} MCQs on {topic}. Format: Q, A/B/C/D, Correct Answer."
        quiz_content = ask_local_llm(prompt)
    return render_template('quiz.html', quiz=quiz_content)

@app.route('/summarize', methods=['GET', 'POST'])
def summarize():
    summary, video_id = None, None
    if request.method == 'POST':
        video_url = request.form.get('video_url', '')
        match = re.search(r"(?:v=|\/)([0-9A-Za-z_-]{11})", video_url)
        if match:
            video_id = match.group(1)
            # Use Gemini's strength: If you have a transcript, send it.
            # For the demo, we will summarize the topic based on URL info
            prompt = f"Summarize the educational content usually found in this video/topic: {video_url}"
            summary = ask_local_llm(prompt)
            session['last_summary'] = summary
    return render_template('summary.html', summary=summary, video_id=video_id)

@app.route('/chat', methods=['GET', 'POST'])
def chat():
    user_name = session.get('user_name', 'Guest')
    if request.method == 'POST':
        user_msg = request.json.get('message') if request.is_json else request.form.get('message')
        image_file = request.files.get('image')
        base64_image = base64.b64encode(image_file.read()).decode('utf-8') if image_file else None
        prompt = f"User: {user_name}. Question: {user_msg}"
        reply = ask_local_llm(prompt, image_data=base64_image)
        return jsonify({"response": reply})
    return render_template('chat.html', user_name=user_name)

if __name__ == '__main__':
    app.run(port=5000)
